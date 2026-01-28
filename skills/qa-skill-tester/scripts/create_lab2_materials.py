from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_lab2_pitch():
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Laboratory Work 2: The QA Skill Tester"
    subtitle.text = "Building the Ultimate Automated Validation System\nLesson 7"
    
    # Slide 2: The Problem
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "The Problem: Agent Entropy"
    tf = body_shape.text_frame
    tf.text = "As agent systems grow, complexity exploits:"
    
    p = tf.add_paragraph()
    p.text = "Manual validation is too slow."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Frontend regressions (broken links, overlays) go unnoticed."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Lack of standardized traceability."
    p.level = 1

    # Slide 3: The Solution
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = "The Solution: QA Meta-Skill"
    
    tf = shapes.placeholders[1].text_frame
    tf.text = "A Skill that Tests Skills:"
    
    p = tf.add_paragraph()
    p.text = "1. Generates Plans (Test Architect / Gemini)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Creates Checklists (XLSX Skill)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Reports Bugs (DOCX Skill)"
    p.level = 1

    p = tf.add_paragraph()
    p.text = "4. Archives Evidence (Artifact Archival Code)"
    p.level = 1
    
    output_path = "lab2_pitch_deck.pptx"
    prs.save(output_path)
    print(f"✅ Pitch Deck created: {output_path}")

if __name__ == "__main__":
    create_lab2_pitch()
